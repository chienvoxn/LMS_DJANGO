from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader

from .exceptions import DocumentParseError


@dataclass(frozen=True)
class ParsedUnit:
    text: str
    page_number: int | None = None
    section_title: str = ""


def _parse_pdf(path):
    try:
        reader = PdfReader(path)
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception as exc:
                raise DocumentParseError(
                    "PDF đang được bảo vệ bằng mật khẩu."
                ) from exc

        units = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                units.append(
                    ParsedUnit(text=text, page_number=page_number)
                )
        return units
    except DocumentParseError:
        raise
    except Exception as exc:
        raise DocumentParseError("Không thể đọc tệp PDF.") from exc


def _parse_txt(path):
    raw = Path(path).read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "cp1258", "latin-1"):
        try:
            text = raw.decode(encoding).strip()
            return [ParsedUnit(text=text)] if text else []
        except UnicodeDecodeError:
            continue
    raise DocumentParseError("Không xác định được encoding của tệp TXT.")


def _parse_docx(path):
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentParseError(
            "Thiếu python-docx. Hãy cài requirements-rag.txt."
        ) from exc

    try:
        document = Document(path)
        paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        return [ParsedUnit(text=text)] if text else []
    except Exception as exc:
        raise DocumentParseError("Không thể đọc tệp DOCX.") from exc


def _parse_pptx(path):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentParseError(
            "Thiếu python-pptx. Hãy cài requirements-rag.txt."
        ) from exc

    try:
        presentation = Presentation(path)
        units = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            parts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    parts.append(text.strip())
            if parts:
                units.append(
                    ParsedUnit(
                        text="\n".join(parts),
                        page_number=slide_number,
                        section_title=f"Slide {slide_number}",
                    )
                )
        return units
    except Exception as exc:
        raise DocumentParseError("Không thể đọc tệp PPTX.") from exc


def parse_document(path, file_type):
    extension = (file_type or Path(path).suffix).lower().lstrip(".")
    parsers = {
        "pdf": _parse_pdf,
        "txt": _parse_txt,
        "docx": _parse_docx,
        "pptx": _parse_pptx,
    }
    parser = parsers.get(extension)
    if parser is None:
        raise DocumentParseError(
            f"Định dạng .{extension} chưa được hỗ trợ."
        )

    units = parser(path)
    if not units or not any(unit.text.strip() for unit in units):
        hint = (
            " Tệp có thể là PDF scan chỉ chứa hình ảnh; "
            "phiên bản này chưa hỗ trợ OCR."
            if extension == "pdf"
            else ""
        )
        raise DocumentParseError(
            "Không trích xuất được văn bản từ tài liệu." + hint
        )
    return units
